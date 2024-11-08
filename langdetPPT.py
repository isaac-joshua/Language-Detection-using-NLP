from pptx import Presentation
from pptx.util import Inches

# Initialize the presentation
presentation = Presentation()

# Slide 1: Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Autocorrect Project Using TextBlob"
subtitle.text = "A Natural Language Processing Application"

# Slide 2: Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.placeholders[1]
title.text = "Introduction"
content.text = (
    "This project aims to create a simple autocorrect tool using TextBlob.\n"
    "TextBlob provides an easy-to-use API for NLP tasks, making it suitable for spelling correction."
)

# Slide 3: Project Objectives
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.placeholders[1]
title.text = "Project Objectives"
content.text = (
    "1. Automate the correction of spelling errors.\n"
    "2. Demonstrate word-wise analysis of corrections.\n"
    "3. Create a simple, extensible framework for text correction."
)

# Slide 4: Technologies and Libraries Used
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.placeholders[1]
title.text = "Technologies and Libraries Used"
content.text = (
    "1. Python: Programming language\n"
    "2. TextBlob: For text processing and autocorrection\n"
    "3. NLTK: Provides linguistic data and models used by TextBlob"
)

# Slide 5: Detailed Implementation
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.placeholders[1]
title.text = "Detailed Implementation"
content.text = (
    "1. Input text is processed using TextBlob.\n"
    "2. The 'correct' method identifies and fixes spelling errors.\n"
    "3. Word-wise analysis shows corrected words."
)

# Slide 6: Code Snippet
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.placeholders[1]
title.text = "Code Snippet"
content.text = (
    "from textblob import TextBlob\n"
    "text = 'I havv a speling errror'\n"
    "blob = TextBlob(text)\n"
    "corrected_text = blob.correct()\n"
    "print(corrected_text)"
)

# Slide 7: Output and Results
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.placeholders[1]
title.text = "Output and Results"
content.text = (
    "Original: 'I havv a speling errror'\n"
    "Corrected: 'I have a spelling error'\n"
    "Word-wise corrections are displayed."
)

# Slide 8: Advantages and Limitations
slide_8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.placeholders[1]
title.text = "Advantages and Limitations"
content.text = (
    "Advantages:\n"
    "- Simple API and quick implementation\n"
    "- Pre-trained models\n\n"
    "Limitations:\n"
    "- Accuracy issues with specialized vocabulary\n"
    "- Performance limitations for large texts"
)

# Slide 9: Future Work
slide_9 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.placeholders[1]
title.text = "Future Work"
content.text = (
    "1. Implement grammar checking.\n"
    "2. Improve context-aware corrections.\n"
    "3. Optimize performance for large datasets."
)

# Slide 10: Conclusion
slide_10 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_10.shapes.title
content = slide_10.placeholders[1]
title.text = "Conclusion"
content.text = (
    "The TextBlob-based autocorrect tool is effective for simple spelling corrections.\n"
    "It can be extended to handle more complex NLP tasks and large-scale applications."
)

# Save the presentation
presentation.save("Autocorrect_Project_Presentation.pptx")
print("Presentation created successfully.")