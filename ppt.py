from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
subtitle1 = slide1.placeholders[1]
title1.text = "Language Detection Script"
subtitle1.text = "An overview of automatic language detection using Python"

# Slide 2: Introduction
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "Introduction"
content2 = slide2.placeholders[1]
content2.text = (
    "This script uses the `langdetect` and `langcodes` libraries to automatically detect the language of a text.\n"
    "Features include:\n"
    "- Language detection using `langdetect`\n"
    "- Conversion of language codes to full names using `langcodes`\n"
    "- Error handling for undetectable languages\n"
)
title2.text_frame.paragraphs[0].font.size = Pt(36)
content2.text_frame.paragraphs[0].font.size = Pt(24)
content2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# Slide 3: Code Explanation - `detect_language` Function
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "detect_language Function"
content3 = slide3.placeholders[1]
content3.text = (
    "This function takes a text as input and returns the detected language.\n"
    "- It first detects the language code (e.g., 'en' for English).\n"
    "- Converts the code to a full language name using `langcodes`.\n"
    "- Handles exceptions with an error message."
)
title3.text_frame.paragraphs[0].font.size = Pt(36)
content3.text_frame.paragraphs[0].font.size = Pt(24)
content3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 51)

# Slide 4: Code Explanation - Main Script
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "Main Script"
content4 = slide4.placeholders[1]
content4.text = (
    "The main script runs multiple example texts to test language detection.\n"
    "- Loops through sample sentences in different languages.\n"
    "- Calls `detect_language` on each sentence and prints the results.\n"
    "- Displays detected language name and code."
)
title4.text_frame.paragraphs[0].font.size = Pt(36)
content4.text_frame.paragraphs[0].font.size = Pt(24)
content4.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 51, 0)

# Slide 5: Conclusion and Usage
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Conclusion & Usage"
content5 = slide5.placeholders[1]
content5.text = (
    "This script can be useful for:\n"
    "- Language-based filtering\n"
    "- Pre-processing text in multilingual applications\n"
    "- Developing language-aware applications\n\n"
    "Run the script with different texts to test its detection capabilities!"
)
title5.text_frame.paragraphs[0].font.size = Pt(36)
content5.text_frame.paragraphs[0].font.size = Pt(24)
content5.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 0, 102)

# Save the presentation
prs.save("Language_Detection_Script_Presentation.pptx")
print("Presentation created: Language_Detection_Script_Presentation.pptx")
